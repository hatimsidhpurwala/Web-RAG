import io
import tempfile
from pathlib import Path
from config.settings import GROQ_API_KEY, TESSERACT_CMD, WHISPER_MODEL

_PDF_EXT   = {".pdf"}
_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}
_AUDIO_EXT = {".wav", ".mp3", ".m4a", ".ogg", ".webm"}
_WORD_EXT  = {".docx"}
_EXCEL_EXT = {".xlsx", ".xls", ".csv"}
_PPT_EXT   = {".pptx", ".ppt"}
_TEXT_EXT  = {".txt", ".md"}

def file_category(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext in _PDF_EXT:    return "pdf"
    if ext in _IMAGE_EXT:  return "image"
    if ext in _AUDIO_EXT:  return "audio"
    if ext in _WORD_EXT:   return "word"
    if ext in _EXCEL_EXT:  return "excel"
    if ext in _PPT_EXT:    return "powerpoint"
    if ext in _TEXT_EXT:   return "text"
    return "unknown"

def extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages).strip()

def extract_image_ocr(raw: bytes) -> str:
    import cv2
    import numpy as np
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    arr = np.frombuffer(raw, dtype=np.uint8)
    img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    denoised = cv2.medianBlur(thresh, 3)
    return pytesseract.image_to_string(denoised, lang="eng").strip()

def extract_word(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_excel(raw: bytes, filename: str) -> str:
    import pandas as pd
    ext = Path(filename).suffix.lower()
    if ext == ".csv":
        df = pd.read_csv(io.BytesIO(raw))
    else:
        df = pd.read_excel(io.BytesIO(raw))
    return df.to_string(index=False)

def extract_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text(raw: bytes) -> str:
    return raw.decode("utf-8", errors="ignore")

def transcribe_audio(raw: bytes, suffix: str) -> str:
    from groq import Groq
    client = Groq(api_key=GROQ_API_KEY)
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    with open(tmp_path, "rb") as f:
        result = client.audio.transcriptions.create(
            file=(f"audio{suffix}", f),
            model=WHISPER_MODEL,
            language="en",
            response_format="text",
        )
    Path(tmp_path).unlink(missing_ok=True)
    return result.strip() if isinstance(result, str) else str(result).strip()

def extract_text_from_file(raw: bytes, filename: str) -> tuple[str, str]:
    """Return (extracted_text, category). Raises ValueError if unsupported."""
    cat = file_category(filename)
    if cat == "pdf":        return extract_pdf(raw), cat
    if cat == "image":      return extract_image_ocr(raw), cat
    if cat == "word":       return extract_word(raw), cat
    if cat == "excel":      return extract_excel(raw, filename), cat
    if cat == "powerpoint": return extract_pptx(raw), cat
    if cat == "text":       return extract_text(raw), cat
    raise ValueError(f"Unsupported file type: {Path(filename).suffix}")
