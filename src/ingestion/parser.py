"""
src/ingestion/parser.py
ALL input parsing: PDF, Word, Excel, PPT, OCR, Audio.
"""

import io
import tempfile
from pathlib import Path

from src.config import WHISPER_MODEL, GROQ_API_KEY, TESSERACT_CMD


def extract_from_pdf(raw: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages).strip()


def extract_from_word(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_from_excel(raw: bytes, filename: str) -> str:
    import pandas as pd
    ext = Path(filename).suffix.lower()
    if ext == ".csv":
        df = pd.read_csv(io.BytesIO(raw))
    else:
        df = pd.read_excel(io.BytesIO(raw))
    return df.to_string(index=False)


def extract_from_ppt(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    texts = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    ]
    return "\n".join(texts)


def extract_from_txt(raw: bytes) -> str:
    return raw.decode("utf-8", errors="ignore")


def extract_from_image(raw: bytes) -> str:
    import cv2
    import numpy as np
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

    arr = np.frombuffer(raw, dtype=np.uint8)
    img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(
        gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
    )
    denoised = cv2.medianBlur(thresh, 3)
    return pytesseract.image_to_string(denoised, lang="eng").strip()


def transcribe_audio(raw: bytes, suffix: str) -> str:
    from groq import Groq
    client = Groq(api_key=GROQ_API_KEY)

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name

    try:
        with open(tmp_path, "rb") as f:
            result = client.audio.transcriptions.create(
                file=(f"audio{suffix}", f),
                model=WHISPER_MODEL,
                language="en",
                response_format="text",
            )
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    return result.strip() if isinstance(result, str) else str(result).strip()


def parse(raw: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return extract_from_pdf(raw)
    if ext == ".docx":
        return extract_from_word(raw)
    if ext in {".xlsx", ".xls", ".csv"}:
        return extract_from_excel(raw, filename)
    if ext in {".pptx", ".ppt"}:
        return extract_from_ppt(raw)
    if ext in {".txt", ".md"}:
        return extract_from_txt(raw)
    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}:
        return extract_from_image(raw)
    if ext in {".wav", ".mp3", ".m4a", ".ogg", ".webm"}:
        return transcribe_audio(raw, ext)

    raise ValueError(f"Unsupported file type: {ext}")
